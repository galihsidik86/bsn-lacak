"""Patch BSN-Lacak-ArtiVisi.pptx supaya akurat dengan sistem yang sudah
dibangun. Pakai python-pptx + edit text run level supaya formatting
(font, warna, bold) preserved. Idempotent — jalankan ulang OK.

Strategi: text-search-replace eksak. Setiap entry di REPLACEMENTS
match satu run, ganti isi, pertahankan style run aslinya."""
from __future__ import annotations
import sys
from pathlib import Path
from pptx import Presentation

SRC = Path(r"C:/Users/Galih Sidik/Downloads/BSN-Lacak-ArtiVisi.pptx")
DST = Path(r"C:/Users/Galih Sidik/Downloads/BSN-Lacak-ArtiVisi-updated.pptx")

# Daftar replacement (old_substring → new). Match substring di run text.
# Kalau old tidak ditemukan, ignore (idempotent).
REPLACEMENTS: list[tuple[str, str]] = [
    # --- Slide 1: Update tanggal dokumen ---
    ("17 Juni 2026", "30 Juni 2026"),

    # --- Slide 3: Solusi modul ---
    # Modul 02 lintas-platform claim
    ("Aplikasi lapangan lintas-platform (Android & iOS) dengan proteksi anti-tampering native.",
     "Aplikasi lapangan: PWA browser + APK Android native (Capacitor) dengan foreground service GPS."),
    # Modul 02 bullets — tambah APK FCM
    ("→  Offline-first via IndexedDB queue",
     "→  Offline-first via localStorage queue"),
    ("→  Geo-Fencing check-in radius 50m",
     "→  Background GPS via foreground service"),
    ("→  In-app camera + watermark server",
     "→  In-app camera + watermark sharp"),
    ("→  Rute optimal nearest-neighbor",
     "→  FCM push notification (Firebase)"),

    # --- Slide 5: Arsitektur ---
    # Hilangkan klaim LDAP yang tidak ada
    ("RBAC 3-tier: ADMIN HQ, SUPERVISOR per-cabang, PETUGAS per-binaan. SSO LDAP + TOTP 2FA wajib.",
     "RBAC 3-tier: ADMIN HQ, SUPERVISOR per-cabang, PETUGAS per-binaan. Username/password + TOTP 2FA + JWT 15m."),
    # HA: ganti ke realita single VPS
    ("HA: PostgreSQL Streaming Replication + dual app server + Nginx LB. SLA 99,9%.",
     "Single-host: Postgres 16 + API Node di Docker Compose, Caddy TLS otomatis Let's Encrypt."),
    # Layer 1 stack
    ("Vite + React 18  ·  TypeScript  ·  Mobile PWA  ·  Capacitor Shell  ·  MapTiler GIS  ·  Web Dashboard  ·  IndexedDB Offline",
     "Vite + React 18 · TypeScript · PWA + Capacitor APK · MapLibre + MapTiler · localStorage queue · Wake Lock API"),
    # Layer 2 endpoint count
    ("API & Services (60+ Endpoints)",
     "API & Services (190+ Endpoints)"),
    ("Node.js 20 + Express  ·  TypeScript + Zod  ·  JWT Auth 15m  ·  LDAP/AD SSO  ·  Prisma ORM  ·  RBAC + Anti-Impersonation  ·  sharp Watermarking",
     "Node.js 22 + Express · TypeScript + Zod · JWT 15m + Refresh 7d · TOTP 2FA · Prisma ORM · RBAC + Branch-scoping · sharp Watermarking"),
    # Layer 3
    ("PostgreSQL 16  ·  WAL Replication  ·  WhatsApp Cloud API  ·  Twilio SMS  ·  SSE Real-time  ·  pino + Prometheus",
     "PostgreSQL 16 · SSE Real-time · FCM Push (Firebase) · Web Push VAPID · BlastGateway pluggable (Twilio/WA-ready) · pino structured log"),
    # Bottom row
    ("🐳 Docker Compose  ·  ⚡ Nginx TLS 1.3  ·  🔄 CI/CD GitHub Actions  ·  🔒 CMMI Level 3  ·  Apache 2.0",
     "🐳 Docker Compose · ⚡ Caddy TLS 1.3 + HTTP/3 · 🔄 CI/CD GitHub Actions · 🔒 RBAC + Audit Immutable · MIT"),

    # --- Slide 6: Domain Model ---
    ("14 ENTITAS DATA", "35+ ENTITAS DATA"),
    ("Database BSN Lacak terdiri dari 14 entitas data yang saling terintegrasi. Berikut adalah 8 entitas utama yang menjadi representasi inti sistem.",
     "Database BSN Lacak terdiri dari 35+ entitas data terintegrasi (Branch, User, Petugas, Nasabah, Angsuran, Kunjungan, Attendance, ChatMessage, PushSubscription, Audit, dll). Berikut 8 entitas inti."),
    # Stat block "14 Entitas Total" → 35+
    # NOTE: ada dua "14" di slide ini — yang stat besar di kanan bawah dengan label "Entitas Total"
    # Kita patch lewat shape-level traversal di bawah karena teks angka "14" terlalu generik untuk plain replace
    ("60+\nREST API Endpoints", "190+\nREST API Endpoints"),

    # --- Slide 7: Dashboard ---
    ("14\nHalaman Dashboard", "20+\nHalaman Dashboard"),

    # --- Slide 8: Mobile ---
    ("8 Halaman Aplikasi Petugas", "Aplikasi Petugas — 5 Tab + Overlay"),
    ("Android (Chrome) · iOS (Safari) · Tanpa App Store",
     "PWA browser (Android/iOS) + APK Android native (Capacitor + FCM) — sideload via GitHub Actions"),
    ("Antrian IndexedDB — data tetap masuk di area blank spot. Sinkronisasi otomatis saat online dengan persistensi tab-kill.",
     "Antrian localStorage + retry queue clientTs — data tetap masuk di blank spot. Sinkron otomatis saat online/focus. Wake Lock jaga layar aktif saat clock-in."),

    # --- Slide 9: Anti-Fraud ---
    # Hilangkan freerasp claim
    ("Deteksi perangkat Root/Jailbreak, aplikasi mock GPS & modifikasi binary via Capacitor + plugin freerasp real-time.",
     "Foto wajib via in-app camera (capture attribute). Backend reject upload non-magic-byte JPEG/PNG. EXIF freshness check tolak foto >1 jam."),
    ("freerasp — Real-time", "Capacitor + Magic Byte"),
    ("Anti-Tampering Native", "Anti-Tampering Sederhana"),

    # --- Slide 10: WhatsApp — soften claim ---
    ("Memanfaatkan WhatsApp Business Cloud API resmi Meta dengan template tersetujui — bukan WA gateway ilegal. BlastGateway dirancang pluggable + fallback SMS Twilio.",
     "BlastGateway pluggable — Twilio SMS sudah terintegrasi, slot WhatsApp Business Cloud API siap aktivasi saat BSN punya Meta Business account. Template santun sesuai etika muamalah."),

    # --- Slide 13: API spec ---
    ("60+ REST ENDPOINTS", "190+ REST ENDPOINTS"),

    # --- Slide 11: Roadmap — masih ada LDAP + freerasp di rincian fase ---
    ("·  Anti-Tampering Capacitor + freerasp (2–3 minggu)",
     "·  In-app camera lock + magic-byte server validation (1 minggu)"),
    ("Spesifikasi teknis, LDAP AD sandbox, tim PIC BSN & ArtiVisi",
     "Spesifikasi teknis, env Firebase + VPS, tim PIC BSN & ArtiVisi"),

    # --- Slide 7: Peta GIS bullet 2 — ganti generic click ke GPS trail + heatmap ---
    ("✓  Klik marker → linimasa kunjungan per petugas per hari",
     "✓  GPS trail historis dengan date picker + heatmap kunjungan"),

    # --- Slide 7: Analitik bullet 3 — arsip JSONL diganti chat + alerts ---
    ("✓  Arsip JSONL.gz harian otomatis untuk kepatuhan OJK MRTI",
     "✓  Chat petugas ↔ supervisor + geofence violation + inactivity alert"),
    ("✓  Supervisor approve/reject + catatan + push notif ke petugas",
     "✓  Supervisor approve/reject + push notif FCM ke petugas"),

    # --- Slide 11 P0: Ganti item ke ✓ DONE per baris (paragraph-level) ---
    ("·  Geo-Fencing 50m server-side (1–2 hari)",
     "✓  Geo-Fencing 50m server-side — DONE"),
    ("·  SSO Active Directory + 2FA TOTP (1 minggu)",
     "✓  TOTP 2FA + JWT rotation 15m — DONE"),
    ("·  In-app camera lock + magic-byte server validation (1 minggu)",
     "✓  In-app camera + magic-byte + EXIF freshness — DONE"),

    # --- Slide 11 P1: 2 item done, 1 tetap pending (WA) ---
    ("·  WhatsApp Business Cloud API resmi (2 minggu)",
     "·  WhatsApp Business Cloud API resmi (tunggu akses Meta)"),
    ("·  Smart Allocation 3-parameter (1 minggu)",
     "✓  Smart Allocation 3-parameter — DONE"),
    ("·  Interactive WA template CTA button (1 minggu)",
     "✓  APK Capacitor + Foreground GPS + FCM — DONE"),
    ("·  Scheduler Pre-Due & Past-Due tier DPD (1 minggu)",
     "✓  Chat petugas ↔ supervisor realtime — DONE"),

    # --- Slide 11 P2: Core Banking realistic, HA update ---
    ("·  Core Banking connector REST/SFTP (2 minggu)",
     "·  Core Banking connector REST/webhook (tunggu spec)"),
    ("·  High Availability Postgres replica (1 minggu)",
     "·  High Availability Postgres replica (saat scale > 1 region)"),
]


def patch_runs(prs) -> int:
    """Walk semua run di semua slide, apply REPLACEMENTS substring."""
    n_changed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                # Gabungkan teks paragraph dulu untuk match multi-run text
                full = "".join(r.text for r in para.runs)
                for old, new in REPLACEMENTS:
                    if old in full:
                        # Strategi: ganti di run pertama, kosongkan sisa.
                        # Ini mempertahankan formatting run pertama untuk
                        # seluruh kalimat (acceptable trade-off vs preserve
                        # multi-run formatting yang lebih rumit).
                        replaced = full.replace(old, new)
                        if para.runs:
                            para.runs[0].text = replaced
                            for r in para.runs[1:]:
                                r.text = ""
                        n_changed += 1
                        full = replaced  # supaya replacement berikutnya
                                          # di paragraf yang sama tetap nyambung
    return n_changed


def patch_slide6_entity_count(prs) -> int:
    """Slide 6 punya 2 stat besar standalone — angka '14' dan '60+' yang
    label-nya di shape terpisah. Plain text-replace tidak bisa karena
    angka tersebut juga muncul di slide lain (kredensial ArtiVisi 40+/60+).
    Kita locate by neighbor-label match."""
    changed = 0
    if len(prs.slides) < 6:
        return 0
    slide = prs.slides[5]  # 0-indexed
    shapes = list(slide.shapes)
    targets = [
        ("14",  "Entitas Total",      "35+"),
        ("60+", "REST API Endpoints", "190+"),
    ]
    for i, shape in enumerate(shapes):
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        for old_num, label, new_num in targets:
            if t != old_num:
                continue
            neighbors = []
            for j in range(max(0, i-2), min(len(shapes), i+3)):
                if j == i: continue
                if shapes[j].has_text_frame:
                    neighbors.append(shapes[j].text_frame.text.strip())
            if any(label in n for n in neighbors):
                shape.text_frame.paragraphs[0].runs[0].text = new_num
                changed += 1
                break
    return changed


def main():
    if not SRC.exists():
        print(f"ERROR: source not found: {SRC}", file=sys.stderr)
        sys.exit(1)
    prs = Presentation(str(SRC))
    n1 = patch_runs(prs)
    n2 = patch_slide6_entity_count(prs)
    prs.save(str(DST))
    print(f"Patched {n1} replacements + {n2} entity-count stat.")
    print(f"Output: {DST}")


if __name__ == "__main__":
    main()
