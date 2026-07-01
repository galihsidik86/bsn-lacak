"""Sisip slide "Spesifikasi Infrastruktur" ke BSN-Lacak-ArtiVisi.pptx
di posisi 12 (setelah Roadmap 14 minggu, sebelum kredensial ArtiVisi).

Strategi: clone slide 6 (Domain Model — punya 8 stat cards + 4 bottom
stats yang cocok untuk konten infra), ganti isi teks dengan spek infra
(4 tier server + laptop supervisor + 3 kategori HP + deployment), lalu
reorder ke posisi 12. Update page number semua slide dari "X / 14" ke
"X / 15".

Sumber data: docs/SPESIFIKASI_INFRASTRUKTUR.md (Section 2-5)."""
from __future__ import annotations
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn

SRC = Path(r"C:/Users/Galih Sidik/BSN/docs/BSN-Lacak-ArtiVisi.pptx")
DST = Path(r"C:/Users/Galih Sidik/BSN/docs/BSN-Lacak-ArtiVisi.pptx")  # overwrite in-place


def duplicate_slide(prs, source_idx: int):
    """Duplicate slide by deep-copying shapes + relationships.
    Key trick: rels di source pakai rId tertentu (mis. rId10 untuk image),
    tapi `get_or_add` di dest assign rId berbeda (mis. rId2). Kalau kita
    tidak remap rId di shape XML, PowerPoint deteksi corruption karena
    XML reference rId10 sementara rels declare rId2. Solusinya: build
    mapping (source_rid → new_rid), lalu rewrite semua rId reference
    di deep-copied XML."""
    import re
    source = prs.slides[source_idx]
    dest = prs.slides.add_slide(source.slide_layout)
    for shape in list(dest.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Bangun mapping rId source → rId dest sambil re-add relationships.
    rid_map: dict[str, str] = {}
    for src_rid, rel in source.part.rels.items():
        if "notesSlide" in rel.reltype:
            continue
        # get_or_add akan return existing rId kalau reltype+target sudah ada,
        # atau assign rId baru. Ambil rId hasilnya untuk mapping.
        # get_or_add return string rId (not relationship object) di
        # python-pptx recent versions.
        new_rid = dest.part.rels.get_or_add(rel.reltype, rel.target_part)
        rid_map[src_rid] = new_rid

    # Deep copy shapes lalu rewrite rId reference di XML pakai mapping.
    RID_ATTR_RE = re.compile(r'((?:r:embed|r:link|r:id)=")(rId\d+)(")')
    for shape in source.shapes:
        new_el = deepcopy(shape._element)
        # Serialize → rewrite → parse ulang. Cara paling aman untuk hit
        # semua atribut namespace r: yang mungkin (r:embed untuk pic,
        # r:link untuk hyperlink, r:id untuk ole/chart, dll).
        from lxml import etree
        xml_str = etree.tostring(new_el, encoding='unicode')
        def sub(m):
            old_rid = m.group(2)
            return m.group(1) + rid_map.get(old_rid, old_rid) + m.group(3)
        new_xml = RID_ATTR_RE.sub(sub, xml_str)
        new_el2 = etree.fromstring(new_xml)
        dest.shapes._spTree.insert_element_before(new_el2, 'p:extLst')

    return dest


def move_slide(prs, from_idx: int, to_idx: int):
    """Reorder slide via manipulasi sldIdLst XML."""
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    el = slides[from_idx]
    sld_id_lst.remove(el)
    # Setelah remove, index target menggeser; sesuaikan
    if to_idx > from_idx:
        to_idx -= 1
    # Insert at target position
    sld_id_lst.insert(to_idx, el)


def set_text_preserve_format(shape, new_text: str):
    """Ganti seluruh teks shape, pertahankan formatting run pertama.
    Kalau new_text ada newline, split ke multiple paragraph."""
    tf = shape.text_frame
    lines = new_text.split("\n")
    # Pertahankan format run pertama paragraph 1 sebagai template
    template_para = tf.paragraphs[0]
    template_run_props = None
    template_pPr = None
    if template_para.runs:
        template_run_props = template_para.runs[0]._r.get_or_add_rPr()
    # Simpen pPr paragraph 1 (alignment, dll)
    template_pPr_el = template_para._pPr
    # Clear semua paragraph existing
    txBody = tf._txBody
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)
    # Rebuild dengan lines baru
    from pptx.oxml.ns import qn as _qn
    from lxml import etree
    for i, line in enumerate(lines):
        p = etree.SubElement(txBody, _qn('a:p'))
        if template_pPr_el is not None:
            p.insert(0, deepcopy(template_pPr_el))
        if line:
            r = etree.SubElement(p, _qn('a:r'))
            if template_run_props is not None:
                r.append(deepcopy(template_run_props))
            t = etree.SubElement(r, _qn('a:t'))
            t.text = line


def patch_new_slide(slide):
    """Ganti konten teks slide clone (dari Domain Model) ke Infra Spec.
    Mapping berdasarkan urutan shape yang muncul di slide 6."""
    # Dictionary lookup: old_text_snippet → new_text
    # Kita gak perlu match exact, cukup startswith unik.
    replacements = {
        # Header
        "Domain Model Sistem":
            "Spesifikasi Infrastruktur Sistem",
        # Source pptx sudah pernah di-update di sesi sebelumnya:
        # "14 ENTITAS DATA" jadi "35+ ENTITAS DATA", dan intro paragraph
        # juga direwrite. Kita match state sekarang, bukan aslinya.
        "35+ ENTITAS DATA":
            "DEPLOYMENT & PERANGKAT",
        "Database BSN Lacak terdiri dari 35+ entitas data terintegrasi (Branch, User, Petugas, Nasabah, Angsuran, Kunjungan, Attendance, ChatMessage, PushSubscription, Audit, dll). Berikut 8 entitas inti.":
            "Tiga komponen perangkat menopang BSN Lacak: server backend, workstation supervisor, dan device petugas lapangan. Deployment on-premise di DC BSN dengan footprint minimal 2 GB VPS untuk pilot, scaling jelas hingga skala enterprise multi-cabang.",

        # 8 Cards (map dari 8 entitas ke 8 komponen infra)
        "Branch":
            "Server — Demo/Pilot",
        "Cabang BSN. Root entity untuk multi-tenancy. Admin HQ melihat semua, Supervisor hanya cabangnya.":
            "2 GB RAM · 1-2 vCPU · 30 GB SSD · Bandwidth 1 TB/bulan. Cocok POC ≤ 10 petugas aktif, 1 cabang, < 500 nasabah.",

        "User":
            "Server — Operasional Kecil",
        "Akun: username, bcrypt hash (cost 12), role, branchId, TOTP, lockout 5x gagal.":
            "4 GB RAM · 2 vCPU · 60 GB SSD · 2 TB/bulan. Sweet spot 1 kantor cabang: 10-50 petugas, 500-2.000 nasabah.",

        "Petugas":
            "Server — Skala Menengah",
        "Field collector: userId (1:1), NIP, HP, lastPosition (lat/lng/at) diperbarui tiap GPS ping.":
            "8 GB RAM · 4 vCPU · 120 GB SSD + storage terpisah · 5 TB/bulan. Multi-cabang 50-200 petugas, 5.000+ nasabah.",

        "Nasabah":
            "Server — Enterprise",
        "Nasabah pembiayaan: NIK, akad, kolektabilitas (1-5), outstanding, koordinat, DPD.":
            "16+ GB RAM · 8+ vCPU · DB terpisah + S3 storage · Private link. 200+ petugas, HA replica, disaster recovery.",

        "Angsuran":
            "Laptop Supervisor",
        "Ledger pembayaran: nasabahId, amount, metode (tunai/transfer), createdAt.":
            "RAM 4 GB min (rec 8 GB) · Chrome/Edge/Firefox terbaru · Layar 14\" Full HD · Internet stabil untuk SSE realtime.",

        "Kunjungan":
            "HP Petugas — PWA",
        "Laporan: GPS, riskScore, riskFlags[], foto ber-watermark, reviewStatus Supervisor.":
            "RAM 2 GB · Android 8+ / iOS 14+ · Chrome 90+ / Safari 14+ · GPS aktif. Cakupan tracking ~70-85% (browser suspend saat layar mati).",

        "Blast":
            "HP Petugas — APK Native",
        "Job komunikasi WA/SMS: template, segmen, scheduledAt, status, terkirim, gagal.":
            "RAM 3+ GB · Android 8+ · Google Play Services WAJIB (FCM) · foreground service GPS. Cakupan tracking ~95-99% walau app tertutup.",

        "AuditLog":
            "Deployment & Ops",
        "Log immutable: actor, action, IP, requestId, meta(JSONB). Tidak dapat dimodifikasi.":
            "Docker Compose · Caddy TLS 1.3 + HTTP/3 · PostgreSQL 16 · Backup harian retensi 30 hari · TLS 1.3 Let's Encrypt auto-renew.",

        # 4 Bottom stats — ganti ke CapEx figures. Angka bawah ini
        # standalone di shape terpisah, tidak akan bentrok dengan
        # "35+" di eyebrow (yang sudah kita replace di atas jadi
        # "DEPLOYMENT & PERANGKAT" — jadi replacement "35+" berikutnya
        # aman hanya match bottom stat).
        "35+":
            "Rp 25 jt",
        "Entitas Total":
            "CapEx 1 Cabang (6 petugas)",

        "190+":
            "Rp 3,6 jt",
        "REST API Endpoints":
            "OpEx Server / Tahun",

        "Prisma":
            "On-Prem",
        "Type-safe ORM":
            "Data Center BSN",

        "Cursor":
            "30 Hari",
        "Pagination Standard":
            "Retensi Backup",

        # Page number
        "6 / 14":
            "12 / 15",  # slide baru jadi halaman 12 dari 15
    }

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        # Match by exact para text (setelah trim). Trim newlines.
        # Cari key yang match awal string.
        for old, new in replacements.items():
            if t == old or t == old.strip():
                set_text_preserve_format(shape, new)
                break


def update_all_page_numbers(prs, new_total: int):
    """Slide-slide asli punya string 'X / 14' — update jadi 'X / <new>'.
    Slide 12 (baru sisipan) sudah punya '12 / 15' dari patch di atas;
    slide 12-14 asli (sekarang jadi slide 13-15) juga perlu update
    baik denominator maupun numerator."""
    # Setelah reorder, urutan slide baru:
    # 1-11: sama seperti asli
    # 12: slide baru (Infra Spec)
    # 13-15: slide asli 12-14 (kredensial ArtiVisi, API spec, Closing)
    #
    # Slide asli 1-11 numerator = index+1 (masih benar), denominator perlu update 14→15
    # Slide 13-15 numerator perlu shift +1
    import re
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in slide.shapes[0].text_frame.paragraphs if False else []: pass
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    txt = run.text
                    # Pattern "N / 14" di posisi standalone di slide.
                    m = re.match(r'^\s*(\d+)\s*/\s*14\s*$', txt)
                    if m:
                        # Update numerator ke i, denominator ke new_total
                        run.text = f"{i} / {new_total}"
                    # Sudah handle di patch_new_slide untuk "6 / 14" → "12 / 15"
                    # tapi kalau numerator sudah ke-set 12/15 by patch, biarkan.
                    # Kalau numerator masih "12 / 15" tapi bukan halaman 12 lagi
                    # (nggak mungkin karena kita loop urut), pass.
                    m2 = re.match(r'^\s*(\d+)\s*/\s*15\s*$', txt)
                    if m2:
                        num = int(m2.group(1))
                        if num != i:
                            run.text = f"{i} / 15"


def main():
    if not SRC.exists():
        print(f"ERROR: source not found: {SRC}", file=sys.stderr)
        sys.exit(1)
    prs = Presentation(str(SRC))
    n_slides = len(prs.slides)
    print(f"Before: {n_slides} slides")

    # Idempotency: kalau sudah ada slide "Spesifikasi Infrastruktur", skip
    for s in prs.slides:
        for shape in s.shapes:
            if shape.has_text_frame and "Spesifikasi Infrastruktur Sistem" in shape.text_frame.text:
                print("SKIP: infra slide already exists")
                return

    # Clone slide 6 (index 5) — Domain Model style
    new_slide = duplicate_slide(prs, 5)
    print(f"Cloned slide 6 → new slide at end (index {len(prs.slides)-1})")

    # Patch content
    patch_new_slide(new_slide)
    print("Content patched")

    # Move new slide from end (index n_slides) to position 12 (index 11)
    move_slide(prs, len(prs.slides) - 1, 11)
    print("Moved to position 12")

    # Update all page numbers X/14 → X/15
    update_all_page_numbers(prs, 15)
    print("Page numbers refreshed")

    prs.save(str(DST))
    print(f"Saved: {DST}")
    print(f"After: {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
